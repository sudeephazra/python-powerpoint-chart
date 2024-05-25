from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Load the presentation
prs = Presentation('samples/python-test.pptx')

# Access the first slide (change index if needed)
slide = prs.slides[0]

# Access the first shape which should be a chart (change index if needed)
chart_shape = None
for shape in slide.shapes:
    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        chart_shape = shape
        break

if chart_shape is None:
    raise ValueError("No chart found in the slide")

# Access the chart
chart = chart_shape.chart
# Define new data
chart_data = CategoryChartData()
chart_data.categories = ['Category 1', 'Category 2', 'Category 3']
chart_data.add_series('Series 1', (30, 40, 30))


chart.replace_data(chart_data)
#
# # Save the presentation with the updated chart data
prs.save('samples/updated_presentation.pptx')
